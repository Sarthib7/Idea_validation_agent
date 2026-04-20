"""Pitch deck and file analysis tool"""
import logging
from typing import Dict, Any
from crewai.tools import tool
import json
import tempfile
import os

logger = logging.getLogger(__name__)


@tool("Pitch Deck Analyzer")
def file_analyzer_tool(file_url: str) -> str:
    """
    Download and analyze a pitch deck file (PDF or PPTX).

    Args:
        file_url: URL to the pitch deck file

    Returns:
        JSON string with deck analysis including content, structure, and quality assessment
    """
    try:
        import httpx

        if not file_url or not (file_url.startswith('http://') or file_url.startswith('https://')):
            return json.dumps({
                "error": "Invalid URL",
                "message": "File URL must start with http:// or https://"
            })

        logger.info(f"Downloading file from: {file_url}")

        # Download file
        try:
            response = httpx.get(file_url, timeout=30.0, follow_redirects=True)
            response.raise_for_status()
            file_content = response.content

            # Detect file type
            content_type = response.headers.get('content-type', '').lower()
            file_ext = None

            if 'pdf' in content_type or file_url.lower().endswith('.pdf'):
                file_ext = '.pdf'
            elif 'presentation' in content_type or 'powerpoint' in content_type or file_url.lower().endswith('.pptx'):
                file_ext = '.pptx'
            else:
                return json.dumps({
                    "error": "Unsupported file type",
                    "message": "Only PDF and PPTX files are supported",
                    "content_type": content_type
                })

            # Save to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
                tmp_file.write(file_content)
                tmp_path = tmp_file.name

            try:
                if file_ext == '.pdf':
                    result = _analyze_pdf(tmp_path)
                else:
                    result = _analyze_pptx(tmp_path)

                result["file_url"] = file_url
                result["file_type"] = file_ext[1:]  # Remove dot

                return json.dumps(result, indent=2)

            finally:
                # Clean up temp file
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)

        except httpx.HTTPError as e:
            return json.dumps({
                "error": "Failed to download file",
                "message": str(e)
            })

    except ImportError as e:
        return json.dumps({
            "error": f"Missing dependency: {e}",
            "message": "Install httpx, pymupdf4llm, and python-pptx"
        })
    except Exception as e:
        logger.error(f"File analyzer error: {e}")
        return json.dumps({"error": str(e)})


def _analyze_pdf(pdf_path: str) -> Dict[str, Any]:
    """Analyze PDF pitch deck"""
    try:
        import pymupdf4llm

        # Extract markdown from PDF
        md_text = pymupdf4llm.to_markdown(pdf_path)

        # Basic analysis
        lines = md_text.split('\n')
        non_empty_lines = [line for line in lines if line.strip()]

        # Try to identify slides (look for markdown headers or page breaks)
        slides = []
        current_slide = []

        for line in non_empty_lines:
            if line.startswith('#'):  # Likely a slide title
                if current_slide:
                    slides.append('\n'.join(current_slide))
                    current_slide = []
            current_slide.append(line)

        if current_slide:
            slides.append('\n'.join(current_slide))

        # Assess deck structure
        assessment = _assess_deck_structure(md_text, len(slides))

        return {
            "status": "success",
            "slide_count": len(slides),
            "content_preview": md_text[:1000],
            "structure_assessment": assessment,
            "extracted_text_length": len(md_text)
        }

    except Exception as e:
        logger.error(f"PDF analysis failed: {e}")
        return {
            "status": "failed",
            "error": f"Could not analyze PDF: {str(e)}"
        }


def _analyze_pptx(pptx_path: str) -> Dict[str, Any]:
    """Analyze PPTX pitch deck"""
    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)

        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slides_content.append({
                "slide_number": i + 1,
                "content": '\n'.join(slide_text)
            })

        # Combine all text
        all_text = '\n\n'.join([s["content"] for s in slides_content])

        # Assess deck structure
        assessment = _assess_deck_structure(all_text, len(slides_content))

        return {
            "status": "success",
            "slide_count": len(slides_content),
            "slides": slides_content[:5],  # First 5 slides
            "content_preview": all_text[:1000],
            "structure_assessment": assessment
        }

    except Exception as e:
        logger.error(f"PPTX analysis failed: {e}")
        return {
            "status": "failed",
            "error": f"Could not analyze PPTX: {str(e)}"
        }


def _assess_deck_structure(full_text: str, slide_count: int) -> Dict[str, Any]:
    """Assess pitch deck structure against Sequoia framework"""

    text_lower = full_text.lower()

    # Check for key sections (Sequoia 10-point framework)
    sections_found = {
        "problem": any(keyword in text_lower for keyword in ['problem', 'pain', 'challenge']),
        "solution": any(keyword in text_lower for keyword in ['solution', 'product', 'how it works']),
        "market": any(keyword in text_lower for keyword in ['market', 'tam', 'opportunity']),
        "business_model": any(keyword in text_lower for keyword in ['business model', 'revenue', 'pricing', 'monetization']),
        "competition": any(keyword in text_lower for keyword in ['competition', 'competitor', 'alternative']),
        "team": any(keyword in text_lower for keyword in ['team', 'founder', 'about us']),
        "traction": any(keyword in text_lower for keyword in ['traction', 'growth', 'metrics', 'users', 'revenue']),
        "financial": any(keyword in text_lower for keyword in ['financial', 'projection', 'forecast'])
    }

    sections_count = sum(sections_found.values())

    # Assess quality
    quality_issues = []
    quality_strengths = []

    if slide_count < 10:
        quality_issues.append("Deck may be too short (typically 10-15 slides for VC pitch)")
    elif slide_count > 20:
        quality_issues.append("Deck may be too long (typically 10-15 slides for VC pitch)")
    else:
        quality_strengths.append("Appropriate slide count")

    if sections_count >= 6:
        quality_strengths.append("Covers most key pitch sections")
    else:
        quality_issues.append(f"Missing key sections (found {sections_count}/8)")

    # Check for specific missing critical sections
    if not sections_found["problem"]:
        quality_issues.append("Missing clear problem statement")
    if not sections_found["solution"]:
        quality_issues.append("Missing solution description")
    if not sections_found["market"]:
        quality_issues.append("Missing market size/opportunity")

    return {
        "sections_found": sections_found,
        "completeness_score": round((sections_count / 8) * 10, 1),
        "quality_issues": quality_issues,
        "quality_strengths": quality_strengths
    }
